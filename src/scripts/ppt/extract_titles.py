import os
from pathlib import Path
from pptx import Presentation
import json
from datetime import datetime

def write_log(log_file, message):
    """Write message to log file with timestamp."""
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    with open(log_file, 'a', encoding='utf-8') as f:
        f.write(f"[{timestamp}] {message}\n")

def emu_to_points(emu):
    """Convert EMU (English Metric Units) to points."""
    return emu / 12700  # 1 point = 12700 EMU

def process_shape(shape, log_file, indent=""):
    """Process a shape and its children if it's a group shape."""
    shape_type = type(shape).__name__
    
    # Get shape position and size in points
    top = emu_to_points(shape.top) if hasattr(shape, 'top') else 0
    height = emu_to_points(shape.height) if hasattr(shape, 'height') else 0
    
    # Log shape details
    position_info = f"top={top:.2f}pt, height={height:.2f}pt"
    write_log(log_file, f"{indent}Type={shape_type}, Position: {position_info}")
    
    text = ""
    if hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            write_log(log_file, f"{indent}  Text: '{text}'")
    
    # If it's a group shape, process its children
    if shape_type == 'GroupShape' and hasattr(shape, 'shapes'):
        write_log(log_file, f"{indent}  Group contains:")
        for child_shape in shape.shapes:
            process_shape(child_shape, log_file, indent + "    ")
    
    return {
        'text': text,
        'top': top,
        'height': height,
        'type': shape_type
    }

def extract_titles_from_ppt(ppt_path, log_file):
    """Extract titles from PowerPoint slides, skipping first and last slides."""
    write_log(log_file, f"\nProcessing: {ppt_path}")
    write_log(log_file, "=" * 80)
    
    prs = Presentation(ppt_path)
    titles = {}
    
    # Skip first and last slides
    content_slides = list(prs.slides)[1:-1]
    write_log(log_file, f"Found {len(content_slides)} content slides (excluding first and last)\n")
    
    for i, slide in enumerate(content_slides, 2):  # Start from 2 since we skipped first slide
        write_log(log_file, f"\nProcessing Slide {i}:")
        write_log(log_file, "-" * 40)
        
        # Debug all shapes
        write_log(log_file, f"Shapes found in slide {i}:")
        title_candidates = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Get raw EMU values for comparison
            top = shape.top if hasattr(shape, 'top') else 0
            height = shape.height if hasattr(shape, 'height') else 0
            
            write_log(log_file, f"\nShape {shape_idx}:")
            text = shape.text.strip() if hasattr(shape, "text") else ""
            
            # Log in points for readability
            write_log(log_file, f"  Type={type(shape).__name__}, Position: top={emu_to_points(top):.2f}pt, height={emu_to_points(height):.2f}pt")
            if text:
                write_log(log_file, f"    Text: '{text}'")
            
            # Consider shapes that:
            # 1. Have text
            # 2. Are in roughly the same position as our known title (around 299079 EMU)
            # 3. Have similar height characteristics (around 916878 EMU)
            if text:
                # Check if shape has similar characteristics to known title shape
                is_at_title_position = 250000 <= top <= 350000  # Around 299079 EMU
                has_title_height = 800000 <= height <= 1000000  # Around 916878 EMU
                
                if is_at_title_position and has_title_height:
                    title_candidates.append({
                        'text': text,
                        'top': top,
                        'height': height
                    })
                    write_log(log_file, f"  Found potential title: '{text}'")
                    write_log(log_file, f"  Position: {emu_to_points(top):.2f}pt, Height: {emu_to_points(height):.2f}pt")
        
        # Process title candidates
        if title_candidates:
            write_log(log_file, "\nTitle candidates found:")
            for idx, candidate in enumerate(title_candidates):
                write_log(log_file, f"  {idx+1}. '{candidate['text']}' (top: {emu_to_points(candidate['top']):.2f}pt)")
            
            # If multiple candidates, choose the one closest to our known good position (299079 EMU)
            title_candidates.sort(key=lambda x: abs(x['top'] - 299079))
            title_text = title_candidates[0]['text']
            titles[i] = title_text
            write_log(log_file, f"\nSelected title: '{title_text}' (position: {emu_to_points(title_candidates[0]['top']):.2f}pt)")
        else:
            write_log(log_file, "\nNo title candidates found in slide!")
    
    return titles

def process_ppt_directory(ppt_dir, log_file):
    """Process all PPT files in directory."""
    ppt_dir = Path(ppt_dir)
    results = {}
    
    for ppt_file in ppt_dir.glob("*.ppt*"):
        try:
            titles = extract_titles_from_ppt(ppt_file, log_file)
            results[ppt_file.name] = {
                "file_path": str(ppt_file),
                "titles": titles
            }
            
            # Print summary
            write_log(log_file, "\nSummary:")
            write_log(log_file, "=" * 40)
            write_log(log_file, f"Total slides with titles: {len(titles)}")
            for slide_num, title in titles.items():
                write_log(log_file, f"Slide {slide_num}: {title}")
            
        except Exception as e:
            write_log(log_file, f"Error processing {ppt_file}: {e}")
            import traceback
            write_log(log_file, traceback.format_exc())
    
    return results

def main():
    # Create log file
    log_file = "data/videos/ppt_extraction.log"
    os.makedirs(os.path.dirname(log_file), exist_ok=True)
    
    # Clear previous log
    with open(log_file, 'w', encoding='utf-8') as f:
        f.write(f"Starting PPT extraction at {datetime.now()}\n")
    
    # Process PPTs from the data/videos/ppts directory
    ppt_dir = "data/videos/ppts"
    results = process_ppt_directory(ppt_dir, log_file)
    
    # Save results
    output_file = "data/videos/ppt_titles.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    write_log(log_file, f"\nResults saved to {output_file}")
    print(f"Process complete. Check {log_file} for detailed log.")

if __name__ == "__main__":
    main() 